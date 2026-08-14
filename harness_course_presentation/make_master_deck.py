import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]

C_DARK_BG   = RGBColor(15, 23, 42)      # Slate 900
C_CARD_BG   = RGBColor(30, 41, 59)      # Slate 800
C_LIGHT_BG  = RGBColor(248, 250, 252)  # Slate 50
C_WHITE     = RGBColor(255, 255, 255)
C_PRIMARY   = RGBColor(15, 23, 42)
C_SECONDARY = RGBColor(51, 65, 85)    # Slate 700
C_MUTED     = RGBColor(148, 163, 184)     # Slate 400
C_ACCENT    = RGBColor(13, 148, 136)     # Teal 600
C_BORDER    = RGBColor(203, 213, 225)    # Slate 300
C_RED       = RGBColor(185, 28, 28)

def make_cover(title, subtitle, schedule):
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK_BG; bg.line.fill.background()

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(4.2), Inches(0.38))
    pill.fill.solid(); pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background()
    p = pill.text_frame.paragraphs[0]
    p.text = '2026 MASTERCLASS: CLAUDE CODE HARNESS'; p.font.size = Pt(10.5); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(8.4), Inches(1.1))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = C_WHITE

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(8.4), Inches(0.6))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = subtitle; p2.font.size = Pt(12.5); p2.font.color.rgb = C_MUTED

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(8.4), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = C_ACCENT; line.line.fill.background()

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.85), Inches(8.4), Inches(2.4))
    card.fill.solid(); card.fill.fore_color.rgb = C_CARD_BG; card.line.color.rgb = C_SECONDARY
    tf = card.text_frame; tf.word_wrap = True
    
    p3 = tf.paragraphs[0]
    p3.text = f"Masterclass Focus: Scaffolding, Guardrails, MCP & Compound Workflows"; p3.font.size = Pt(11); p3.font.bold = True; p3.font.color.rgb = C_ACCENT; p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf.add_paragraph()
    p4.text = schedule; p4.font.size = Pt(10.5); p4.font.color.rgb = C_WHITE; p4.alignment = PP_ALIGN.CENTER
    
    p5 = tf.add_paragraph()
    p5.text = "GitHub Repository: https://github.com/kenhuangus/packt-harness"; p5.font.size = Pt(10.5); p5.font.bold = True; p5.font.color.rgb = C_ACCENT; p5.alignment = PP_ALIGN.CENTER

    p6 = tf.add_paragraph()
    p6.text = "Course Website: https://kenhuangus.github.io/packt-harness/"; p6.font.size = Pt(10.5); p6.font.bold = True; p6.font.color.rgb = C_WHITE; p6.alignment = PP_ALIGN.CENTER
    return slide

def make_section(mod_num, title, timeline, objectives, idx):
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK_BG; bg.line.fill.background()

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(3.8), Inches(0.4))
    pill.fill.solid(); pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background()
    p = pill.text_frame.paragraphs[0]
    p.text = f'MODULE {mod_num}  |  {timeline}'; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_WHITE

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(8.4), Inches(2.4))
    card.fill.solid(); card.fill.fore_color.rgb = C_CARD_BG; card.line.color.rgb = C_SECONDARY
    tf = card.text_frame; tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = 'KEY MODULE OBJECTIVES & GITHUB IMPLEMENTATION'; p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = C_ACCENT
    for obj in objectives:
        pi = tf.add_paragraph()
        pi.text = f'• {obj}'; pi.font.size = Pt(11); pi.font.color.rgb = C_WHITE

    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(5.1), Inches(0.4), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = C_ACCENT; badge.line.fill.background()
    pb = badge.text_frame.paragraphs[0]
    pb.text = str(idx); pb.font.size = Pt(11); pb.font.bold = True; pb.font.color.rgb = C_WHITE; pb.alignment = PP_ALIGN.CENTER
    return slide

def make_cards(title, timeline, cards, idx):
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_LIGHT_BG; bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.4), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = C_PRIMARY

    if timeline:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(0.45), Inches(2.2), Inches(0.38))
        pill.fill.solid(); pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background()
        pp = pill.text_frame.paragraphs[0]
        pp.text = timeline; pp.font.size = Pt(10); pp.font.bold = True; pp.font.color.rgb = C_WHITE; pp.alignment = PP_ALIGN.CENTER

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(8.8), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = C_BORDER; line.line.fill.background()

    n = len(cards)
    if n == 1:
        c = cards[0]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.25), Inches(8.8), Inches(3.7))
        card.fill.solid(); card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
        tf = card.text_frame; tf.word_wrap = True
        pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(13); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
        for bullet in c['bullets']:
            pb = tf.add_paragraph()
            pb.text = f'• {bullet}'; pb.font.size = Pt(11); pb.font.color.rgb = C_SECONDARY
    elif n == 2:
        for i, c in enumerate(cards):
            x = Inches(0.6) if i == 0 else Inches(5.2)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.25), Inches(4.2), Inches(3.7))
            card.fill.solid(); card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
            tf = card.text_frame; tf.word_wrap = True
            pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(12); pt.font.bold = True
            pt.font.color.rgb = C_RED if 'Failure' in c['title'] or 'Warning' in c['title'] or 'Anti' in c['title'] else C_ACCENT
            for bullet in c['bullets']:
                pb = tf.add_paragraph()
                pb.text = f'• {bullet}'; pb.font.size = Pt(10.5); pb.font.color.rgb = C_SECONDARY
    elif n == 3:
        for i, c in enumerate(cards):
            y = Inches(1.25 + i * 1.22)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(8.8), Inches(1.05))
            card.fill.solid(); card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
            tf = card.text_frame; tf.word_wrap = True
            pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(12); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
            for bullet in c['bullets']:
                pb = tf.add_paragraph()
                pb.text = f'• {bullet}'; pb.font.size = Pt(10.5); pb.font.color.rgb = C_PRIMARY
    elif n >= 4:
        for i, c in enumerate(cards[:4]):
            x = Inches(0.6) if i % 2 == 0 else Inches(5.2)
            y = Inches(1.25) if i < 2 else Inches(2.75)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.2), Inches(1.35))
            card.fill.solid(); card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
            tf = card.text_frame; tf.word_wrap = True
            pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(11.5); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
            for bullet in c['bullets']:
                pb = tf.add_paragraph()
                pb.text = f'• {bullet}'; pb.font.size = Pt(9.5); pb.font.color.rgb = C_PRIMARY

    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(5.1), Inches(0.4), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = C_ACCENT; badge.line.fill.background()
    pb = badge.text_frame.paragraphs[0]
    pb.text = str(idx); pb.font.size = Pt(11); pb.font.bold = True; pb.font.color.rgb = C_WHITE; pb.alignment = PP_ALIGN.CENTER
    return slide

slides_cnt = 0

# Slide 1: Cover
make_cover(
    'Build Reliable Claude Code Workflows with Guardrails and Tests',
    'Harness Engineering Masterclass: Scaffolding, Guardrails, MCP & Compound Workflows',
    'Schedule: 09:00 AM - 11:30 AM  |  Production AI Coding Agent Systems Architecture'
)
slides_cnt += 1

# Slide 2: Master Course Map
make_cards('COURSE MASTER MAP & GITHUB REPOSITORY', '09:00 AM - 11:30 AM', [
    {'title': 'PART 1: FOUNDATIONS & CONTROL (09:00 - 10:20)', 'bullets': [
        'Mod 1 (09:00-09:15): Why Harness Engineering & Failure Modes',
        'Mod 2 (09:15-09:30): Core Harness Stack & 5 System Pillars',
        'Mod 3 (09:30-09:50): Spec-Driven Development (SDD) & Schemas',
        'Mod 4 (09:50-10:10): Guardrails, AST Checks & Deterministic Hooks',
        'Break (10:10-10:20): Mid-Morning Break & Permission Gateways Q&A'
    ]},
    {'title': 'PART 2: RELIABILITY & TEAMS (10:20 - 11:30)', 'bullets': [
        'Mod 5 (10:20-10:35): Tests as Reliability Layer & TDA Loop',
        'Mod 6 (10:35-10:50): Skills, Plugins & Model Context Protocol (MCP)',
        'Mod 7 (10:50-11:05): Compound Engineering & Multi-Agent Swarms',
        'Mod 8 (11:05-11:15): 5-Step Practical Workflow SOP Pattern',
        'Closing (11:15-11:30): Core Principles, Audit Checklist & GitHub Code'
    ]}
], slides_cnt + 1)
slides_cnt += 1

# --- MODULE 1 (Slides 3-8) ---
make_section(1, 'Why Harness Engineering', '09:00 AM - 09:15 AM', [
    'Understand why AI model capability does not equal production reliability.',
    'Examine the 98% Harness Rule in 2026 agentic systems engineering.',
    'Analyze concrete failure modes: context drift, infinite loops, & unverified edits.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_01_why_harness_engineering/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Speed Without Structure Creates Risk: Harness Thesis', None, [
    {'title': 'The Un-Harnessed Agent Problem', 'bullets': [
        'AI coding agents move fast, but raw model capability without scaffolding creates risk.',
        'Symptom: Impressive single-turn generation, followed by context amnesia & unverified edits.',
        'Risk: Infinite execution loops, path traversal exploits, and broken builds.'
    ]},
    {'title': 'The Deterministic Harness Solution', 'bullets': [
        'The difference between a useful coding agent and an unreliable one isn\'t the model alone — it\'s the harness around it.',
        'Harness engineering builds surrounding systems that make agent behavior predictable, testable, and constrained.',
        'Optimizes for trust, correctness, and production readiness rather than raw speed.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Model vs System Harness: The 2026 Paradigm', None, [
    {'title': 'The Model (Probabilistic Reasoner)', 'bullets': [
        'Generates code proposals based on next-token probability.',
        'Susceptible to hallucinations, context amnesia, and command retries.',
        'Requires external scaffolding to verify logic.'
    ]},
    {'title': 'The Harness (Deterministic Scaffolding)', 'bullets': [
        'Enforces execution boundaries, AST linting, and path sandboxing.',
        'Intercepts infinite retry loops automatically.',
        'Provides standardized LLM client interface for enterprise model providers.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 1: Context Decay & Amnesia', None, [
    {'title': 'Symptoms of Context Failure', 'bullets': [
        'Agent forgets initial constraints after 8+ conversation turns.',
        'Mutates unrelated files or reverses prior architecture decisions.',
        'Mitigation: Persistent memory files (CLAUDE.md / AGENTS.md).'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 2: Infinite Execution Loops', None, [
    {'title': 'Symptoms of Execution Traps', 'bullets': [
        'Agent retries failing compiler command endlessly with minor tweaks.',
        'Burns tokens without making architectural progress.',
        'Mitigation: Harness loop detector in harness_vs_model_demo.py.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 3: Unsanitized Mutations & Bad Tools', None, [
    {'title': 'Symptoms of Dangerous Tool Actions', 'bullets': [
        'Agent attempts rm -rf / or unauthorized shell commands.',
        'Overwrites production configuration files or drops database tables.',
        'Mitigation: Pre-action regex shell interceptors & path isolation.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 2 (Slides 9-14) ---
make_section(2, 'Core Harness Stack & Architecture', '09:15 AM - 09:30 AM', [
    'Deconstruct the 5 core pillars of a production AI coding harness.',
    'Implement persistent memory standards (CLAUDE.md vs AGENTS.md).',
    'Build least-privilege permission models and pre/post hooks.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_02_core_harness_stack/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('System Architecture: The 5 Harness Pillars', None, [
    {'title': 'Instructions & Repo Conventions', 'bullets': [
        'CLAUDE.md / AGENTS.md specifying architecture & rules.',
        'Cascading lookup and symlink interoperability.'
    ]},
    {'title': 'Scoped Tools & Permissions', 'bullets': [
        'Least-privilege tool execution model.',
        'Path sandboxing blocking path traversal exploits.'
    ]},
    {'title': 'Hooks & Policy Engine', 'bullets': [
        'Pre-action security checks.',
        'Post-action AST linting and secret scanning.'
    ]},
    {'title': 'Testing & Tracing', 'bullets': [
        'Automated pytest runner feeding tracebacks back.',
        'JSONL audit event logging for full observability.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 1: Memory Files (CLAUDE.md vs AGENTS.md)', None, [
    {'title': 'CLAUDE.md (Anthropic Ecosystem)', 'bullets': [
        'Dedicated standing brief for Claude Code agent sessions.',
        'Supports cascading subdirectory discovery.'
    ]},
    {'title': 'AGENTS.md (Universal Open Standard)', 'bullets': [
        'Cross-tool standard for Cursor, Copilot, and Aider.',
        'Symlink pattern: ln -s AGENTS.md CLAUDE.md'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 2: Scoped Tool & Permission Models', None, [
    {'title': 'Principle of Least Privilege Execution', 'bullets': [
        'Grant minimal required tools (read_file, write_file, run_test).',
        'Validate target path resides within workspace root.',
        'Codebase: core_harness_stack.py validate_tool_permission()'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 3: Deterministic Hooks & Policy Engine', None, [
    {'title': 'Pre-Execution & Post-Execution Interceptors', 'bullets': [
        'Pre-Hooks: Inspect tool call arguments before execution.',
        'Post-Hooks: Trigger AST static analysis and secret scanners on edits.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillars 4 & 5: Context Token Budgeting & Tracing', None, [
    {'title': 'Context Window Compacting & Budgeting', 'bullets': [
        'Allocate token budget: 20% Memory, 20% Spec, 50% Edit, 10% Output.',
        'Truncate long compiler outputs while retaining core traceback.'
    ]},
    {'title': 'JSONL Observability & Audit Tracing', 'bullets': [
        'Log structured JSON entries for every tool call and error.',
        'Audit trail logged to events.jsonl for complete telemetry.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 3 (Slides 15-21) ---
make_section(3, 'Spec-Driven Development (SDD)', '09:30 AM - 09:50 AM', [
    'Master Spec-Driven Development to eliminate prompt ambiguity.',
    'Draft machine-verifiable acceptance criteria and scope boundaries.',
    'Implement upfront spec verification and anti-hallucination checks.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_03_spec_driven_development/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Why Prompts Fail & Specs Succeed', None, [
    {'title': 'The Flaws of Natural Language Prompts', 'bullets': [
        'Prompts are ambiguous, subjective, and context-dependent.',
        'Lead to feature creep and unrequested refactoring.'
    ]},
    {'title': 'The Power of Executable Specifications', 'bullets': [
        'Defines clear contracts, allowed files, and non-goals.',
        'Codebase: SPEC.md & spec_driven_verifier.py'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Anatomy of a Production SPEC.md', None, [
    {'title': 'Required SPEC.md Structure', 'bullets': [
        'Executive Summary & Problem Context',
        'Allowed Modification Scope (Explicit file list)',
        'Explicit Non-Goals (Forbidden actions)',
        'Input/Output Data Schemas (JSON/TypeScript)',
        'Machine-Verifiable Acceptance Criteria'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Defining Scope Boundaries & Non-Goals', None, [
    {'title': 'Hard Boundaries to Prevent Scope Creep', 'bullets': [
        'Coding agents tend to refactor adjacent files needlessly.',
        'SPEC.md lists explicitly allowed files (e.g. auth_validator.py).',
        'Verifier rejects edits to database.py or config files.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Drafting Machine-Verifiable Acceptance Criteria', None, [
    {'title': 'Vague vs Executable Criteria', 'bullets': [
        'Vague: Make the API fast and handle errors gracefully.',
        'Executable: validate_jwt(token) returns {"valid": True} for valid tokens.',
        'Executable: 100% test coverage in tests/test_auth.py.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('How Claude Code & AGY Process Specs', None, [
    {'title': 'Agent Execution Lifecycle under SDD', 'bullets': [
        'Agent reads SPEC.md at session initialization.',
        'Agent generates implementation plan matching SPEC criteria.',
        'Spec verifier audits diff against allowed scope prior to merge.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Spec Verification & Anti-Hallucination Checks', None, [
    {'title': 'Upfront Consistency Verification', 'bullets': [
        'Run spec validation pass prior to code generation.',
        'Ensure proposed diff satisfies non-goals and schema constraints.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 4 (Slides 22-29) ---
make_section(4, 'Guardrails and Deterministic Hooks', '09:50 AM - 10:10 AM', [
    'Enforce 4-layer defense-in-depth security boundaries.',
    'Build pre-action shell interceptors and post-action AST linters.',
    'Implement path isolation sandboxing and secret key scanning.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_04_guardrails_and_hooks/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('The 4-Layer Control Architecture', None, [
    {'title': 'Layer 1: System Rules (Prompt)', 'bullets': [
        'Soft guidelines defined in CLAUDE.md / AGENTS.md.'
    ]},
    {'title': 'Layer 2: Tool Schemas (JSON Schema)', 'bullets': [
        'Strict argument type & payload shape enforcement.'
    ]},
    {'title': 'Layer 3: Interceptors & Hooks (Process)', 'bullets': [
        'Pre-action regex filters & post-action AST static analysis.'
    ]},
    {'title': 'Layer 4: OS Sandboxing (System)', 'bullets': [
        'Path traversal blocks & workspace root chroot isolation.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Claude Code Native Hooks & Shell Interception', None, [
    {'title': 'Claude Code Event Hooks (PreToolUse & PostToolUse)', 'bullets': [
        'Native events include PreToolUse and PostToolUse; matcher selects tool names such as Bash.',
        'Explicitly block --dangerously-skip-permissions CLI flag.'
    ]},
    {'title': 'Inspecting Shell Commands Before Execution', 'bullets': [
        'Regex blocklist targeting rm -rf, sudo, chmod 777, git push --force.',
        'Codebase: guardrails_engine.py intercept_shell_command()'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Post-Action Checks: AST Parsers & Linters', None, [
    {'title': 'AST Static Analysis & Secret Scanning', 'bullets': [
        'Runs automatically whenever agent writes or updates a file.',
        'Executes ast.parse() to catch syntax & structural errors.',
        'Scans for hardcoded secret API keys (sk-proj...).'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Filesystem Sandboxing & Path Isolation', None, [
    {'title': 'Workspace Boundary Enforcement', 'bullets': [
        'Bind agent file operations strictly to project working directory.',
        'Block path traversal exploits (../../etc/passwd or C:\\Windows).'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Admin Escalation & Interactive Approvals', None, [
    {'title': 'Risk-Tiered Permission Matrix', 'bullets': [
        'Auto-approve low-risk tools: read_file, list_dir, grep.',
        'Log medium-risk tools: write_file, run_test.',
        'Require user confirmation modal for critical state mutations.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('SWE-bench Evaluation & Guardrail Verification', None, [
    {'title': 'Benchmarking Agent Control Harnesses', 'bullets': [
        'SWE-bench research proves test-only evaluation is insufficient.',
        'AST static analysis checks architectural & rationale integrity.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Production Safety & Audit Matrix', None, [
    {'title': 'Verification Checklist', 'bullets': [
        '[PASS] Are all shell commands logged to audit JSONL with timestamps?',
        '[PASS] Is secret scanning enforced by post-action checks and git pre-commit checks?',
        '[PASS] Is path traversal sandboxing active for all file write tools?'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 5 (Slides 30-31) ---
make_section(5, 'Mid-Morning Break & Permission Gateways', '10:10 AM - 10:20 AM', [
    'Take a 10-minute break to refresh.',
    'Participate in open Q&A on Session 1: Harness Stack, SDD & Guardrails.',
    'Review Risk-Tiered Permission Gateway implementation.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_05_break_and_qna/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Open Q&A & Group Discussion', None, [
    {'title': 'Discussion Prompts', 'bullets': [
        'How are you currently managing agent permissions in your organization?',
        'What failure modes have caused the most friction in your workflows?',
        'How can spec-driven development accelerate your team velocity?'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 6 (Slides 32-37) ---
make_section(6, 'Tests as Reliability Layer', '10:20 AM - 10:35 AM', [
    'Architect Test-Driven Agent (TDA) execution loops.',
    'Leverage multi-tier testing (Unit, Integration, E2E) for validation.',
    'Extract automated tracebacks and register anti-regression tests.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_06_tests_as_reliability_layer/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Test-Driven Agent (TDA) Loop Architecture', None, [
    {'title': 'The TDA Execution Loop', 'bullets': [
        'Write/Identify failing unit test based on SPEC.md.',
        'Agent generates implementation code to satisfy test.',
        'Run test runner; extract traceback on failure.',
        'Feed traceback back into agent repair prompt.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Multi-Tier Testing for Coding Agents', None, [
    {'title': 'Testing Hierarchy', 'bullets': [
        'Unit Tests: Fast isolation tests verifying individual functions.',
        'Integration Tests: Verifies multi-module interactions & APIs.',
        'E2E Tests: Full system workflow verification.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Automated Error Traceback Feedback Loops', None, [
    {'title': 'Zero-Touch Error Capture', 'bullets': [
        'Eliminates manual copy-pasting of compiler/test errors by developer.',
        'Harness captures pytest/jest stdout and stderr automatically.',
        'Passes formatted traceback directly into LLM repair prompt.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Turning Agent Failures into Permanent Tests', None, [
    {'title': 'Anti-Regression Safeguards', 'bullets': [
        'Every agent hallucination or bug must yield a new regression test.',
        'Add test case to project regression test suite permanently.',
        'Codebase: tda_reliability_pipeline.py register_anti_regression_test()'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Test Reliability Scorecard & Metrics', None, [
    {'title': 'Key Performance Indicators (KPIs)', 'bullets': [
        'First-Pass Pass Rate: % of agent edits passing tests on turn 1.',
        'Mean Iterations to Repair: Avg attempts to fix failing test (< 3).',
        'Regression Rate: % of edits breaking existing tests.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 7 (Slides 38-45) ---
make_section(7, 'Skills, Plugins, and MCP Tools', '10:35 AM - 10:50 AM', [
    'Master the Agent Skills open standard (SKILL.md).',
    'Package skills, agents, hooks, MCP/LSP servers, and monitors into modular plugins.',
    'Implement Model Context Protocol (MCP) servers over stdio & Streamable HTTP.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_07_skills_plugins_mcp/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Agent Skills Open Standard (SKILL.md)', None, [
    {'title': 'Folder-Based Skill Encapsulation', 'bullets': [
        'Self-contained folder containing SKILL.md and reference assets.',
        'YAML frontmatter defines metadata; description drives skill discovery.',
        'Encapsulates domain expertise and repeatable workflows.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Structuring Skill Instructions & Resources', None, [
    {'title': 'Skill Directory Anatomy', 'bullets': [
        'SKILL.md: Frontmatter metadata + operational instructions.',
        'references/: Deep documentation, schema specs, and design rules.',
        'scripts/: Executable helper scripts and automation tools.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Plugin Bundles: Packaging Capabilities', None, [
    {'title': 'Deployable Extension Packages', 'bullets': [
        'Bundles Skills, Agents, Hooks, MCP/LSP Servers, and Monitors into deployable units.',
        'Configured via .claude-plugin/plugin.json manifest.',
        'Provides modular domain extensions for enterprise repos.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Model Context Protocol (MCP) Foundation', None, [
    {'title': 'Open Standard by Anthropic', 'bullets': [
        'Open standard created by Anthropic for AI-to-system connectivity.',
        'JSON-RPC 2.0 protocol separating message format from transport.',
        'Exposes Tools, Resources, and Prompts to LLMs.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('MCP Transports: stdio vs Streamable HTTP', None, [
    {'title': 'Transport Layer Comparison', 'bullets': [
        'stdio Transport: Standard for local desktop tools and CLI extensions.',
        'Streamable HTTP: Current remote transport; HTTP+SSE is its legacy predecessor.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('MCP Server Implementation Example', None, [
    {'title': 'MCP Python SDK 2.x', 'bullets': [
        'from mcp.server.mcpserver import MCPServer',
        'mcp = MCPServer("harness-tools")',
        '@mcp.tool() defines callable tools; mcp.run() starts stdio.',
        'Codebase: mcp_server_demo.py & mcp_client_runner.py'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('MCP Governance & Security Controls', None, [
    {'title': 'Enterprise Governance', 'bullets': [
        'OAuth2 & API Key authentication for remote MCP servers.',
        'Granular tool-level permission gating (read vs write actions).'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 8 (Slides 46-53) ---
make_section(8, 'Compound Engineering & Agent Teams', '10:50 AM - 11:05 AM', [
    'Design multi-agent architectures for complex software projects.',
    'Implement specialized agent roles: Planner, Implementer, & Reviewer.',
    'Isolate agent workspaces using Git Worktrees.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_08_compound_engineering/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('Multi-Agent Workflows vs Single Agents', None, [
    {'title': 'Cognitive Specialization', 'bullets': [
        'Single agents suffer cognitive overload as file count grows.',
        'Context window bloat reduces reasoning accuracy.',
        'Multi-agent teams partition responsibilities cleanly.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: Planner Agent', None, [
    {'title': 'Architect & Requirement Analyst', 'bullets': [
        'Read-only access to codebase; analyzes SPEC.md requirements.',
        'Drafts high-level architecture & file modification plan.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: Implementer Agent', None, [
    {'title': 'Constrained Coder Agent', 'bullets': [
        'Receives focused subtask assignment from Planner.',
        'Executes file edits within strict workspace boundary or Git Worktree.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: Reviewer Agent', None, [
    {'title': 'Independent Quality Auditor', 'bullets': [
        'Independent agent inspecting Implementer output.',
        'Runs static analysis, AST linters, and secret scanners.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Subagent Context, Tools & Worktree Isolation', None, [
    {'title': 'Native Subagent Isolation Controls', 'bullets': [
        'isolation: worktree gives the subagent its own Git worktree.',
        'Each subagent gets its own context window and focused system prompt.',
        'tools and disallowedTools constrain access; disallowedTools applies first.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Git Worktree Commands & Context Hygiene', None, [
    {'title': 'Subagent Prompt Isolation & Clean Context', 'bullets': [
        'Planner passes ONLY focused sub-spec segment to Implementer.',
        'git worktree add -b feature-agent ./agent-worktree main',
        'Codebase: multi_agent_team_simulator.py'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Building a Recursive Self-Improvement Loop', None, [
    {'title': 'Telemetry-Driven Learning', 'bullets': [
        'Collect execution logs, tracebacks, and developer edits in telemetry.jsonl.',
        'Automatically analyze recurring failure patterns.',
        'Feed learnings back into CLAUDE.md / AGENTS.md rules.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 9 (Slides 54-58) ---
make_section(9, 'Practical Workflow Pattern', '11:05 AM - 11:15 AM', [
    'Master the 5-step standard operating procedure for AI coding agents.',
    'Combine Spec First, Sandboxing, Hooks, Tests, & Human Review.',
    'Execute production SOP pipeline end-to-end.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_09_practical_workflow_pattern/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('The 5-Step Practical Workflow Overview', None, [
    {'title': 'Standard Operating Procedure (SOP)', 'bullets': [
        'Spec First (Draft & validate SPEC.md)',
        'Constrained Execution (Sandboxed agent runner)',
        'Deterministic Checks (Pre/post hooks & AST linters)',
        'Test Verification (Automated pytest suite)',
        'Human Review (Developer sanity check & PR merge)'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phases 1 & 2: Spec First & Constrained Exec', None, [
    {'title': 'Alignment & Sandboxing', 'bullets': [
        'Align human & AI on exact deliverables before writing code.',
        'Run agent with minimal tool permissions & strict file scopes.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phases 3 & 4: Deterministic Checks & Tests', None, [
    {'title': 'Automated Verification Gating', 'bullets': [
        'Automated AST parsers, linters, & security rule gates.',
        'Pass unit/integration test suite; feed back any errors.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phase 5: Human Review & Sanity Auditing', None, [
    {'title': 'Developer Final Approval', 'bullets': [
        'Developer reviews clean, pre-validated PR diff.',
        'Verifies business logic alignment and edge case handling.'
    ]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 10 (Slides 59-63) ---
make_section(10, 'Closing Principles & Production Audit', '11:15 AM - 11:30 AM', [
    'Review the 4 core principles of Harness Engineering.',
    'Evaluate your team\'s readiness using the Production Readiness Audit Suite.',
    'Access full course implementations, multi-provider LLM client & repository.',
    'Lab demo: https://github.com/kenhuangus/packt-harness/blob/main/course_implementation/module_10_closing_and_principles/README.md'
], slides_cnt + 1); slides_cnt += 1

make_cards('The 4 Core Principles of Harness Engineering', None, [
    {'title': 'Foundational Pillars', 'bullets': [
        'Predictability Over Randomness: Standardize CLAUDE.md / AGENTS.md.',
        'Reduce Ambiguity: Use executable SPEC.md instead of prompts.',
        'Automate Checks: Replace human vigilance with deterministic hooks & tests.',
        'Optimize for Trust: Prioritize auditability and correctness.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Production Readiness Audit Checklist', None, [
    {'title': 'Audit Criteria (100% Score)', 'bullets': [
        '[PASS] Is CLAUDE.md / AGENTS.md configured in project root?',
        '[PASS] Are pre-execution hooks active for dangerous shell commands?',
        '[PASS] Is automated test runner integrated into agent repair loop?',
        '[PASS] Are MCP tool permissions scoped appropriately?',
        '[PASS] Is multi-agent role division active for complex subtasks?'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Harness System Verification & Observability', None, [
    {'title': 'Production Readiness Scorecard', 'bullets': [
        'Memory Files (CLAUDE.md / AGENTS.md) configured & symlinked.',
        'Path isolation sandboxing blocking unauthorized file access.',
        'Pre-action regex shell interceptors & post-action AST linters active.',
        'Automated pytest traceback extractor feeding into repair prompts.'
    ]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Final Q&A, Wrap-Up & Course Links', None, [
    {'title': 'Course Resources & Links', 'bullets': [
        'Course Website: https://kenhuangus.github.io/packt-harness/',
        'GitHub Repository: https://github.com/kenhuangus/packt-harness',
        'Includes 10 Module Code Implementations, LLM Client, & Test Suite.',
        'Interactive Web Dashboard: course_implementation/dashboard/index.html'
    ]}
], slides_cnt + 1); slides_cnt += 1

out_dir = r"C:\Users\kenhu\packt\harness\harness_course_presentation\slides\output"
os.makedirs(out_dir, exist_ok=True)
master_path = os.path.join(out_dir, "harness_engineering_master_course_updated.pptx")
prs.save(master_path)

print(f"SUCCESSFULLY RE-GENERATED MASTER DECK ({slides_cnt} Slides) AT:\n{master_path}")
