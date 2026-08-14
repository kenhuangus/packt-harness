import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]

C_DARK_BG = RGBColor(15, 23, 42) # Slate 900
C_CARD_BG = RGBColor(30, 41, 59) # Slate 800
C_LIGHT_BG = RGBColor(248, 250, 252) # Slate 50
C_WHITE = RGBColor(255, 255, 255)
C_PRIMARY = RGBColor(15, 23, 42)
C_SECONDARY = RGBColor(51, 65, 85) # Slate 700
C_MUTED = RGBColor(148, 163, 184) # Slate 400
C_ACCENT = RGBColor(13, 148, 136) # Teal 600
C_BORDER = RGBColor(203, 213, 225) # Slate 300
C_RED = RGBColor(185, 28, 28)

def make_cover(title, subtitle, schedule):
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg.fill.solid; bg.fill.fore_color.rgb = C_DARK_BG; bg.line.fill.background

pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(3.4), Inches(0.4))
pill.fill.solid; pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background
p = pill.text_frame.paragraphs[0]
p.text = 'MASTERCLASS WORKSHOP'; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.4))
p = txBox.text_frame.paragraphs[0]
p.text = title; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = C_WHITE

txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8.4), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = subtitle; p2.font.size = Pt(15); p2.font.color.rgb = C_MUTED

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.04))
line.fill.solid; line.fill.fore_color.rgb = C_ACCENT; line.line.fill.background

card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(8.4), Inches(1.0))
card.fill.solid; card.fill.fore_color.rgb = C_CARD_BG; card.line.color.rgb = C_SECONDARY
p3 = card.text_frame.paragraphs[0]
p3.text = schedule; p3.font.size = Pt(12); p3.font.color.rgb = C_MUTED; p3.alignment = PP_ALIGN.CENTER
return slide

def make_section(mod_num, title, timeline, objectives, idx):
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg.fill.solid; bg.fill.fore_color.rgb = C_DARK_BG; bg.line.fill.background

pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(3.8), Inches(0.4))
pill.fill.solid; pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background
p = pill.text_frame.paragraphs[0]
p.text = f'MODULE {mod_num} | {timeline}'; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(1.0))
p = txBox.text_frame.paragraphs[0]
p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_WHITE

card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(8.4), Inches(2.4))
card.fill.solid; card.fill.fore_color.rgb = C_CARD_BG; card.line.color.rgb = C_SECONDARY
tf = card.text_frame; tf.word_wrap = True
p2 = tf.paragraphs[0]
p2.text = '🎯 KEY MODULE OBJECTIVES'; p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = C_ACCENT
for obj in objectives:
pi = tf.add_paragraph
pi.text = f'• {obj}'; pi.font.size = Pt(12); pi.font.color.rgb = C_WHITE

badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(5.1), Inches(0.4), Inches(0.4))
badge.fill.solid; badge.fill.fore_color.rgb = C_ACCENT; badge.line.fill.background
pb = badge.text_frame.paragraphs[0]
pb.text = str(idx); pb.font.size = Pt(11); pb.font.bold = True; pb.font.color.rgb = C_WHITE; pb.alignment = PP_ALIGN.CENTER
return slide

def make_cards(title, timeline, cards, idx):
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg.fill.solid; bg.fill.fore_color.rgb = C_LIGHT_BG; bg.line.fill.background

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(6.5), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = title; p.font.size = Pt(19); p.font.bold = True; p.font.color.rgb = C_PRIMARY

pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(0.45), Inches(2.2), Inches(0.38))
pill.fill.solid; pill.fill.fore_color.rgb = C_ACCENT; pill.line.fill.background
pp = pill.text_frame.paragraphs[0]
pp.text = timeline; pp.font.size = Pt(10); pp.font.bold = True; pp.font.color.rgb = C_WHITE; pp.alignment = PP_ALIGN.CENTER

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(8.8), Inches(0.02))
line.fill.solid; line.fill.fore_color.rgb = C_BORDER; line.line.fill.background

n = len(cards)
if n == 1:
c = cards[0]
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.25), Inches(8.8), Inches(3.7))
card.fill.solid; card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
tf = card.text_frame; tf.word_wrap = True
pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(14); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
for bullet in c['bullets']:
pb = tf.add_paragraph
pb.text = f'• {bullet}'; pb.font.size = Pt(12); pb.font.color.rgb = C_SECONDARY
elif n == 2:
for i, c in enumerate(cards):
x = Inches(0.6) if i == 0 else Inches(5.2)
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.25), Inches(4.2), Inches(3.7))
card.fill.solid; card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
tf = card.text_frame; tf.word_wrap = True
pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(13); pt.font.bold = True
pt.font.color.rgb = C_RED if 'Failure' in c['title'] or 'Warning' in c['title'] or 'Anti' in c['title'] else C_ACCENT
for bullet in c['bullets']:
pb = tf.add_paragraph
pb.text = f'• {bullet}'; pb.font.size = Pt(11); pb.font.color.rgb = C_SECONDARY
elif n == 3:
for i, c in enumerate(cards):
y = Inches(1.25 + i * 1.22)
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(8.8), Inches(1.05))
card.fill.solid; card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
tf = card.text_frame; tf.word_wrap = True
pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(13); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
for bullet in c['bullets']:
pb = tf.add_paragraph
pb.text = f'• {bullet}'; pb.font.size = Pt(11); pb.font.color.rgb = C_PRIMARY
elif n >= 4:
for i, c in enumerate(cards[:4]):
x = Inches(0.6) if i % 2 == 0 else Inches(5.2)
y = Inches(1.25) if i < 2 else Inches(2.75)
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.2), Inches(1.35))
card.fill.solid; card.fill.fore_color.rgb = C_WHITE; card.line.color.rgb = C_BORDER
tf = card.text_frame; tf.word_wrap = True
pt = tf.paragraphs[0]; pt.text = c['title']; pt.font.size = Pt(12); pt.font.bold = True; pt.font.color.rgb = C_ACCENT
for bullet in c['bullets']:
pb = tf.add_paragraph
pb.text = f'• {bullet}'; pb.font.size = Pt(10); pb.font.color.rgb = C_PRIMARY

badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(5.1), Inches(0.4), Inches(0.4))
badge.fill.solid; badge.fill.fore_color.rgb = C_ACCENT; badge.line.fill.background
pb = badge.text_frame.paragraphs[0]
pb.text = str(idx); pb.font.size = Pt(11); pb.font.bold = True; pb.font.color.rgb = C_WHITE; pb.alignment = PP_ALIGN.CENTER
return slide

slides_cnt = 0

# Slide 1: Cover
make_cover('Harness Engineering for AI Coding Agents', '60+ Slide Comprehensive Architectural Blueprint & Production Guide', '🕒 | 🎯 Claude Code, AGY, MCP & Safety Guardrails | 💡 62 Slides Deck')
slides_cnt += 1

# Slide 2: Master Course Map
make_cards('COURSE MASTER MAP & MAP', '', [
{'title': 'PART 1: FOUNDATIONS & CONTROL ', 'bullets': [
'Mod 1: Why Harness Engineering & Agent Failure Modes',
'Mod 2: Core Harness Stack & 5 System Pillars',
'Mod 3: Spec-Driven Development (SDD) & Machine Specs',
'Mod 4: Guardrails & Deterministic Hooks (Pre/Post Controls)',
'Break: Break & Checkpoint Q&A'
]},
{'title': 'PART 2: RELIABILITY & TEAMS ', 'bullets': [
'Mod 5: Tests as Reliability Layer & Test-Driven Agents',
'Mod 6: Skills, Plugins & Model Context Protocol (MCP)',
'Mod 7: Compound Engineering & Multi-Agent Teams',
'Mod 8: 5-Step Practical Workflow Pattern',
'Closing: Core Principles, Checklist & Final Q&A'
]}
], slides_cnt + 1)
slides_cnt += 1

# --- MODULE 1 (Slides 3-8) ---
make_section(1, 'Why Harness Engineering', '', [
'Define Harness Engineering and why model intelligence alone is insufficient.',
'Understand the 98% Harness Rule in production AI agent systems.',
'Identify common agent failure modes: context drift, infinite loops, & unverified edits.'
], slides_cnt + 1); slides_cnt += 1

make_cards('The Paradigm Shift: Model vs System Harness', '', [
{'title': '🧠 The Model (Probabilistic Engine)', 'bullets': [
'Generates code suggestions based on token probability.',
'Prone to hallucination, context drift, and non-deterministic logic.',
'Lacks native awareness of filesystem boundaries or enterprise policies.'
]},
{'title': '🛡️ The Harness (Deterministic System)', 'bullets': [
'Surrounds model with memory, tools, guardrails, & test verification.',
'Enforces hard policy checks and pre/post action security hooks.',
'Converts probabilistic output into deterministic production commits.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('The 98% Harness Rule', '', [
{'title': '⚙️ Why Systems Engineering Beats Prompting', 'bullets': [
'98% of agent reliability comes from scaffolding, not raw model weights.',
'Prompt tweaks fail to stop infinite loops or file corruption.',
'Harness engineering provides predictable execution environments.',
'Shifts focus from asking nicely to enforcing mechanical constraints.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 1: Context Drift & Amnesia', '', [
{'title': '⚠️ Symptoms of Context Failure', 'bullets': [
'Agent forgets original user constraints after 10+ turns.',
'Introduces conflicting code changes or undoes previous work.',
'Fails to read project conventions in CLAUDE.md or AGENTS.md.'
]},
{'title': '🛡️ Harness Mitigation Strategy', 'bullets': [
'Inject persistent context files into system prompt on every step.',
'Enforce compact execution trajectories & subagent handoffs.',
'Maintain explicit SPEC.md references throughout conversation.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 2: Infinite Execution Loops', '', [
{'title': '⚠️ Symptoms of Execution Loops', 'bullets': [
'Agent retries broken command endlessly with minor syntax tweaks.',
'Consumes token quota while making zero forward progress.',
'Gets trapped editing wrong file or misinterpreting compiler error.'
]},
{'title': '🛡️ Harness Mitigation Strategy', 'bullets': [
'Implement deterministic loop-detection hooks in runner.',
'Set max step limits per tool invocation.',
'Force agent to halt and request human review when stuck.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Failure Mode 3: Unsanitized Edits & Bad Tools', '', [
{'title': '⚠️ Symptoms of Unsafe Tool Calls', 'bullets': [
'Agent attempts rm -rf or executes unapproved shell commands.',
'Modifies global state or edits files outside workspace root.',
'Invokes APIs with hallucinated parameters or incorrect formats.'
]},
{'title': '🛡️ Harness Mitigation Strategy', 'bullets': [
'Strict pre-execution command whitelist & path sandboxing.',
'Schema validation on all tool arguments before invocation.',
'Interactive permission prompts for high-risk operations.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 2 (Slides 9-14) ---
make_section(2, 'Core Harness Stack & Architecture', '', [
'Examine the 5 fundamental pillars of a production agent harness.',
'Configure project memory files (CLAUDE.md, AGENTS.md) and conventions.',
'Implement scoped tool permissions, pre/post hooks, and event tracing.'
], slides_cnt + 1); slides_cnt += 1

make_cards('System Architecture: The 5 Harness Pillars', '', [
{'title': '1. Instructions & Conventions', 'bullets': ['CLAUDE.md, AGENTS.md specifying style, architecture & rules.']},
{'title': '2. Scoped Tools & Permissions', 'bullets': ['Restricted terminal, filesystem, and network tool interfaces.']},
{'title': '3. Hooks & Policy Engine', 'bullets': ['Pre-action security checks & post-action AST validation.']},
{'title': '4. Automated Test Loop', 'bullets': ['Continuous test execution feeding tracebacks back to agent.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 1: Instructions & Repo Conventions', '', [
{'title': '📄 Project Memory Files (CLAUDE.md / AGENTS.md)', 'bullets': [
'Defines build commands, test commands, & formatting guidelines.',
'Provides architecture overviews and module boundaries.',
'Must be kept concise to prevent context bloat.',
'Agents automatically read these files on session startup.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 2: Scoped Tools & Permission Models', '', [
{'title': '🔒 Principle of Least Privilege for Agents', 'bullets': [
'Grant only tool capabilities required for specific task.',
'Isolate file modifications strictly to workspace directory.',
'Disable arbitrary web access or require domain whitelist.',
'Use explicit user approval triggers for destructive commands.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillar 3: Deterministic Hooks & Policy Engine', '', [
{'title': '⚡ Pre-Tool & Post-Tool Interceptors', 'bullets': [
'Pre-Hooks: Inspect CommandLine strings before shell execution.',
'Post-Hooks: Run linters (eslint, ruff) immediately after file edits.',
'Block forbidden imports, API key leaks, or dangerous syntax.',
'Fail fast before code reaches git stage or main branch.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pillars 4 & 5: Automated Testing & Observability', '', [
{'title': '🧪 Automated Testing Loop', 'bullets': [
'Run pytest/jest automatically after every multi-file edit.',
'Supply precise error tracebacks to agent for self-correction.'
]},
{'title': '📊 Observability & Event Tracing', 'bullets': [
'Log all tool calls, token usage, and subagent handoffs in JSONL.',
'Enable post-mortem analysis of agent failure trajectories.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 3 (Slides 15-21) ---
make_section(3, 'Spec-Driven Development (SDD)', '', [
'Learn how to convert loose requirements into machine-executable specs.',
'Define explicit scope boundaries, non-goals, and schema definitions.',
'Guide Claude Code and AGY behavior with automated acceptance criteria.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Fundamentals of Spec-Driven Development', '', [
{'title': '🎯 Why Prompts are Insufficient for Complex Features', 'bullets': [
'Natural language prompts are inherently ambiguous.',
'Agents make assumptions when specifications are incomplete.',
'SDD forces explicit upfront design before code generation.',
'Provides a single source of truth for both Human and AI.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Anatomy of an Executable SPEC.md', '', [
{'title': '📝 Required Sections in a SPEC.md', 'bullets': [
'1. Objective & Background Context',
'2. Functional Requirements & User Stories',
'3. Scope Boundaries & Explicit Non-Goals',
'4. Data Schemas & API Contracts',
'5. Machine-Verifiable Acceptance Criteria'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Defining Scope Boundaries & Non-Goals', '', [
{'title': '🚫 Preventing Scope Creep in Agents', 'bullets': [
'Agents tend to over-engineer solutions or refactor unrelated code.',
'Explicitly state NON-GOALS (e.g. Do NOT modify existing database schemas).',
'Set hard file modification lists (e.g. Edits limited to src/auth/*).',
'Reduces hallucination by narrowing search space.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Drafting Machine-Verifiable Acceptance Criteria', '', [
{'title': '✅ Measurable Criteria vs Vague Goals', 'bullets': [
'Vague: Make user login fast and secure.',
'Verifiable: POST /api/login must return 200 OK with JWT token under 200ms.',
'Verifiable: Unit test coverage in tests/auth_test.py must be > 90%.',
'Allows agent to self-evaluate completion objectively.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('How Claude Code & AGY Consume Specs', '', [
{'title': '🤖 Integration into Agent Workflows', 'bullets': [
'Agent reads SPEC.md at session initiation.',
'Breaks spec into step-by-step implementation plan.',
'Validates each completed subtask against acceptance criteria.',
'Halts execution if code violates spec constraints.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Spec Verification & Anti-Hallucination', '', [
{'title': '🛡️ Spec Audit Checkpoints', 'bullets': [
'Run spec consistency check before implementation.',
'Detect conflicting requirements early.',
'Ensure API parameter types match existing codebase.',
'Guarantees alignment before modifying codebase.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 4 (Slides 22-28) ---
make_section(4, 'Guardrails & Deterministic Hooks', '', [
'Build multi-layered security guardrails for AI coding agents.',
'Implement pre-action command filtering and post-action diff auditing.',
'Enforce filesystem sandboxes, network egress rules, and escalation gates.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Defense-in-Depth Control Architecture', '', [
{'title': '🛡️ Layered Guardrail Design', 'bullets': [
'Layer 1: System Prompt Constraints (Soft Guardrails)',
'Layer 2: Tool Parameter Validation & Schemas (Medium Guardrails)',
'Layer 3: Pre/Post Execution Hooks (Hard Guardrails)',
'Layer 4: OS Sandbox & Container Isolation (Environment Guardrails)'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Pre-Action Hooks: Intercepting Commands', '', [
{'title': '⚡ Intercepting Execution BEFORE it Happens', 'bullets': [
'Evaluates shell command strings prior to terminal execution.',
'Blocks dangerous commands: rm -rf, sudo, chmod, git push --force.',
'Verifies target URLs against allowed egress whitelist.',
'Returns immediate policy violation error to agent.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Post-Action Checks: AST Parsers & Linters', '', [
{'title': '🔍 Auditing Code AFTER Agent Modification', 'bullets': [
'Triggers automatically when agent writes or edits a file.',
'Runs AST static analysis to catch syntax & type errors.',
'Executes secret scanners to prevent API key check-ins.',
'Reverts illegal file changes automatically.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Filesystem Sandboxing & Boundary Isolation', '', [
{'title': '📁 Restricting File Access', 'bullets': [
'Bind agent file operations strictly to current working directory.',
'Prevent path traversal attacks (../../etc/passwd).',
'Protect sensitive files (.env, .git, id_rsa) from reads/writes.',
'Isolate scratch files in temporary build folders.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Admin Escalation & Permission Approval Gates', '', [
{'title': '🚨 Interactive Escalation Triggers', 'bullets': [
'Low-risk actions (read file, run test) auto-approved.',
'Medium-risk actions (edit file, npm install) logged.',
'High-risk actions (git push, db drop) require explicit user click.',
'Maintains developer control over critical operations.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Production Compliance & Auditability Matrix', '', [
{'title': '📋 Compliance & Safety Audit Trail', 'bullets': [
'Log all executed shell commands with cryptographic hashes.',
'Store full prompt/response transcripts for security review.',
'Ensure compliance with SOC2, ISO27001, and internal AI policies.',
'Enable instant rollbacks of agent-generated commits.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 5 (Slides 29-30) ---
make_section(5, ' Break & Open Q&A', '', [
'Take a break to refresh.',
'Engage in open Q&A on Session 1 concepts: Harness Stack, SDD & Hooks.',
'Share real-world agent failure experiences and mitigation strategies.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Open Q&A & Group Discussion', '', [
{'title': '💬 Discussion Questions for Participants', 'bullets': [
'1. What failure modes have you encountered with coding agents?',
'2. How do you currently handle security permissions for LLM tools?',
'3. What are the challenges in adopting Spec-Driven Development?'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 6 (Slides 31-37) ---
make_section(6, 'Tests as Reliability Layer', '', [
'Implement Test-Driven Agent (TDA) workflows.',
'Use multi-tiered test suites (Unit, Integration, E2E) to validate agent output.',
'Convert past agent failures into permanent regression test benchmarks.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Test-Driven Agent (TDA) Architecture', '', [
{'title': '🧪 The TDA Execution Cycle', 'bullets': [
'Step 1: Write failing unit test based on SPEC.md requirements.',
'Step 2: Agent generates implementation code to pass test.',
'Step 3: Run test suite automatically; feed output to agent.',
'Step 4: Refactor until all tests pass cleanly.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Multi-Tier Testing for AI Agents', '', [
{'title': '1. Unit Tests', 'bullets': ['Validate isolated function logic & edge cases. Fast feedback.']},
{'title': '2. Integration Tests', 'bullets': ['Verify interactions between components & DB/API modules.']},
{'title': '3. E2E Regression Tests', 'bullets': ['Ensure full application flows work without breaking.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Automated Error Traceback Feedback Loops', '', [
{'title': '🔄 Converting Stack Traces into Fix Prompts', 'bullets': [
'Never force user to manually copy-paste terminal errors.',
'Harness automatically captures pytest/jest failure output.',
'Formats exact line numbers and tracebacks into next prompt turn.',
'Enables autonomous debugging and self-correction.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Turning Failures into Permanent Tests', '', [
{'title': '🛡️ The Anti-Regression Pipeline', 'bullets': [
'Every time an agent introduces a bug, write a regression test.',
'Add test case to project regression test suite.',
'Prevents future agent runs from re-introducing same bug.',
'Continuously hardens codebase against agent hallucination.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Mitigating Flaky Tests in Sandboxed Runs', '', [
{'title': '⚠️ Managing Test Non-Determinism', 'bullets': [
'Isolate test runs from external network dependencies.',
'Use mock servers for API integration tests.',
'Set explicit timeouts on test execution commands.',
'Distinguish between agent code bugs and infrastructure flakiness.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Test Reliability Scorecard & Metrics', '', [
{'title': '📊 Key Metrics for Agent Test Suites', 'bullets': [
'First-Pass Test Success Rate (target > 70%).',
'Average Repair Iterations per failing test (target < 3).',
'Code Coverage of agent-generated pull requests (target > 85%).',
'Zero Unhandled Exception Rate.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 7 (Slides 38-44) ---
make_section(7, 'Skills, Plugins, and MCP Tools', '', [
'Explore the Agent Skills open standard (SKILL.md).',
'Bundle skills, agents, hooks, MCP/LSP servers, and monitors into modular plugins.',
'Integrate Model Context Protocol (MCP) tools for safe enterprise connectivity.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Agent Skills Open Standard (SKILL.md)', '', [
{'title': '📦 What is an Agent Skill?', 'bullets': [
'A portable, folder-based package containing SKILL.md.',
'Provides agent with specialized domain knowledge & workflows.',
'Loaded dynamically when user prompt matches skill description.',
'Avoids context bloat by reading skill details on-demand.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Structuring Skill Instructions & Resources', '', [
{'title': '📂 Skill Folder Architecture', 'bullets': [
'SKILL.md: YAML metadata; description drives discovery, followed by instructions.',
'references/: Deep documentation, API specifications, & design rules.',
'scripts/: Executable helper scripts & automation tools.',
'assets/ & templates/: Standardized code/file templates.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Plugin Bundles: Packaging Agent Capabilities', '', [
{'title': '🔌 Modular Plugin Architecture', 'bullets': [
'Combines Skills, Agents, Hooks, MCP/LSP Servers, and Monitors into one unit.',
'Configured via .claude-plugin/plugin.json manifest.',
'Enables easy sharing of agent capabilities across engineering teams.',
'Supports versioning and centralized security management.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Model Context Protocol (MCP) Foundation', '', [
{'title': '🌐 Standardized Protocol for Agent Tools', 'bullets': [
'Open standard created by Anthropic for AI-to-system integration.',
'Decouples agent logic from underlying tool implementations.',
'Uses JSON-RPC 2.0 over stdio or Streamable HTTP; HTTP+SSE is legacy.',
'Supports Resources, Tools, and Prompts capabilities.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Connecting Agents to Enterprise APIs & DBs', '', [
{'title': '🔗 Real-World MCP Use Cases', 'bullets': [
'Query internal Postgres/Snowflake databases securely.',
'Fetch GitHub issues, PR diffs, and CI build logs.',
'Interact with AWS/GCP cloud management APIs.',
'Trigger Jenkins pipelines and monitoring alerts.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('MCP Governance, Authentication & Security', '', [
{'title': '🔐 Securing MCP Integrations', 'bullets': [
'OAuth2 & API Key authentication for MCP server connections.',
'Granular tool-level permission gating (read vs write).',
'Rate-limiting and payload size enforcement.',
'Complete audit logging of all MCP request/response payloads.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 8 (Slides 45-52) ---
make_section(8, 'Compound Engineering & Agent Teams', '', [
'Design multi-agent architectures for complex software projects.',
'Implement specialized agent roles: Planner, Implementer, & Reviewer.',
'Build a recursive self-improvement loop for continuous agent optimization.'
], slides_cnt + 1); slides_cnt += 1

make_cards('Multi-Agent Workflows vs Single Agents', '', [
{'title': '👥 Why Single Agents Break on Large Tasks', 'bullets': [
'Single agents suffer context overload when handling complex features.',
'Cognitive load increases exponentially with file count.',
'Multi-agent teams divide work into specialized sub-tasks.',
'Increases task success rate from 40% to over 90%.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: The Planner Agent', '', [
{'title': '🧠 The Architect Subagent', 'bullets': [
'Read-only access to codebase; analyzes requirements.',
'Drafts high-level architecture & file modification plan.',
'Breaks work into discrete, independent subtasks.',
'Does NOT write production code directly.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: The Implementer Agent', '', [
{'title': '⚙️ The Coder Subagent', 'bullets': [
'Receives targeted subtask from Planner.',
'Executes code edits within strict workspace boundary.',
'Runs local unit tests to verify implementation.',
'Reports completed diffs back to team lead.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Role Specialization: The Reviewer Agent', '', [
{'title': '🔍 The Quality & Security Auditor', 'bullets': [
'Independent agent inspecting Implementer output.',
'Runs static analysis, linting, and security audit.',
'Validates code against original SPEC.md acceptance criteria.',
'Approves PR or requests specific code fixes.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Handoff Protocols & Context Passing', '', [
{'title': '📬 Clean Inter-Agent Communication', 'bullets': [
'Pass structured JSON artifacts between subagents.',
'Avoid passing raw conversation histories to keep context light.',
'Include explicit task objectives and file paths in handoff.',
'Enforce sync/async message delivery protocols.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Subagent Context, Tools & Worktree Isolation', '', [
{'title': '1. Context Isolation', 'bullets': ['Each subagent gets its own context window and focused system prompt.']},
{'title': '2. Tool Constraints', 'bullets': ['tools and disallowedTools constrain access; disallowedTools applies first.']},
{'title': '3. Worktree Isolation', 'bullets': ['isolation: worktree gives the subagent its own Git worktree.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Building a Recursive Self-Improvement Loop', '', [
{'title': '🔄 Continuous Agent Hardening', 'bullets': [
'Capture execution telemetry and human correction data.',
'Automatically identify common prompt failure patterns.',
'Update project CLAUDE.md and skill instructions automatically.',
'System gets smarter and more reliable with every project.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 9 (Slides 53-57) ---
make_section(9, 'Practical Workflow Pattern', '', [
'Master the 5-step operational workflow for AI coding agents.',
'Combine Spec First, Constrained Execution, Hooks, Tests, & Human Review.',
'Apply pattern to real-world engineering team workflows.'
], slides_cnt + 1); slides_cnt += 1

make_cards('The 5-Step Practical Workflow Overview', '', [
{'title': '🔄 The Standard Operating Procedure', 'bullets': [
'Step 1: Spec First (Draft & validate SPEC.md)',
'Step 2: Constrained Execution (Sandboxed agent runner)',
'Step 3: Deterministic Checks (Pre/post hooks & linter)',
'Step 4: Test Verification (Automated test suite run)',
'Step 5: Human Review (Final sanity check & PR merge)'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phases 1 & 2: Spec First & Constrained Exec', '', [
{'title': '1. Spec First', 'bullets': ['Align human & AI on exact deliverables before writing code.']},
{'title': '2. Constrained Execution', 'bullets': ['Run agent with minimal tool permissions & strict file scopes.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phases 3 & 4: Deterministic Checks & Tests', '', [
{'title': '3. Deterministic Checks', 'bullets': ['Automated AST parsers, linters, & security rule gates.']},
{'title': '4. Test Verification', 'bullets': ['Pass unit/integration test suite; feed back any errors.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Phase 5: Human Review & Sanity Auditing', '', [
{'title': '👀 The Human-in-the-Loop Safeguard', 'bullets': [
'Developer reviews clean, pre-validated PR diff.',
'Verifies business logic alignment and edge case handling.',
'Merges code into main branch with high confidence.',
'Achieves 10x developer velocity without sacrificing quality.'
]}
], slides_cnt + 1); slides_cnt += 1

# --- MODULE 10 (Slides 58-62) ---
make_section(10, 'Closing Principles & Final Q&A', '', [
'Review the 4 core principles of Harness Engineering.',
'Evaluate your team\'s readiness using the Production Checklist.',
'Participate in final Q&A and wrap-up discussion.'
], slides_cnt + 1); slides_cnt += 1

make_cards('The 4 Core Principles of Harness Engineering', '', [
{'title': '🎯 Predictability Over Randomness', 'bullets': ['Standardize environment & memory files.']},
{'title': '🔍 Reduce Ambiguity', 'bullets': ['Use executable specs instead of natural language prompts.']},
{'title': '⚡ Automate Checks', 'bullets': ['Replace human vigilance with deterministic hooks & test suites.']},
{'title': '🛡️ Optimize for Trust', 'bullets': ['Prioritize auditability and safety over raw speed.']}
], slides_cnt + 1); slides_cnt += 1

make_cards('Production Readiness Checklist', '', [
{'title': '📋 Agent Harness Audit Criteria', 'bullets': [
'✅ Is CLAUDE.md / AGENTS.md configured in project root?',
'✅ Are pre-execution hooks active for dangerous shell commands?',
'✅ Is automated test runner integrated into agent loop?',
'✅ Are MCP tool permissions scoped appropriately?',
'✅ Is multi-agent role division used for complex tasks?'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('The Future of Agentic Software Engineering', '', [
{'title': '🚀 What Lies Ahead in 2026+', 'bullets': [
'Autonomous self-healing software harnesses.',
'Standardized MCP tool marketplaces for enterprise infrastructure.',
'Formal verification of AI-generated code diffs.',
'Developers evolving from coders to Systems Harness Architects.'
]}
], slides_cnt + 1); slides_cnt += 1

make_cards('Final Q&A, Wrap-Up & Workshop Takeaways', '', [
{'title': '💬 Thank You for Participating!', 'bullets': [
'Course Artifacts: All presentation slides & example harness scripts.',
'Skill Package: Global PPT skill installed and ready for use.',
'Open Q&A: Final questions, team discussion & next steps.'
]}
], slides_cnt + 1); slides_cnt += 1

out_path = r'C:\Users\kenhu\packt\harness\harness_course_presentation\slides\output\harness_engineering_62_slides.pptx'
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f'TOTAL SLIDES GENERATED: {slides_cnt}')
print(f'FILE SAVED TO: {out_path}')
